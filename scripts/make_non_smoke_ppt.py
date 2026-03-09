from pathlib import Path
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
FIG_DIR = ROOT / 'figures'
OUT = FIG_DIR / 'fullrun_figures_summary.pptx'

images = sorted([
    p for p in FIG_DIR.glob('*.png')
    if 'smoke' not in p.name.lower()
])

prs = Presentation()

# Cover
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = '论文复现实验结果简报'
subtitle = slide.placeholders[1]
subtitle.text = (
    f'日期：{date.today().isoformat()}\n'
    '姓名：______________\n'
    '内容：使用 figures/ 中全量（非冒烟）图片'
)

# Method brief
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = '说明'
body = slide.shapes.placeholders[1].text_frame
body.clear()
for i, t in enumerate([
    '本简报仅使用全量运行结果图（文件名不含 smoke）。',
    '每张图单独一页，保留原图，附一行简要解读。',
    '重点关注 Simulation 1 与 ETFDR Figure 1 扩展结果。',
]):
    p = body.paragraphs[0] if i == 0 else body.add_paragraph()
    p.text = t
    p.level = 0
    p.font.size = Pt(22)


def infer_note(name: str) -> str:
    n = name.lower()
    if 'selection_rate' in n:
        return '解读：该图展示不同方法在样本规模变化下的正确选模率趋势。'
    if 'nonrejected_count' in n:
        return '解读：该图展示不可拒绝模型数量的变化，用于衡量置信集合宽度。'
    if 'figure1_repro' in n:
        return '解读：复现 ETFDR Figure 1 的主趋势，用于观察误差与变量选择风险的关系。'
    if 'cv_cvc_fdr_mse' in n:
        return '解读：比较 CV 与 CVC 在 FDR-MSE 维度上的差异。'
    if 'ncv1_overlay' in n:
        return '解读：在同一框架下对比 CV/CVC/NCV 的 λ 路径表现。'
    return '解读：该图用于补充验证方法对比结果。'

for img in images:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.5))
    tf = title_box.text_frame
    tf.text = img.name
    tf.paragraphs[0].font.size = Pt(20)

    # image sizing
    max_w, max_h = Inches(12.2), Inches(5.8)
    left, top = Inches(0.55), Inches(0.8)
    pic = slide.shapes.add_picture(str(img), left, top)
    scale = min(max_w / pic.width, max_h / pic.height)
    pic.width = int(pic.width * scale)
    pic.height = int(pic.height * scale)
    pic.left = int((prs.slide_width - pic.width) / 2)
    pic.top = int(Inches(0.8) + (max_h - pic.height) / 2)

    # note
    note_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.55), Inches(12.0), Inches(0.7))
    nt = note_box.text_frame
    nt.text = infer_note(img.name)
    nt.paragraphs[0].font.size = Pt(16)

# Closing
slide = prs.slides.add_slide(prs.slide_layouts[5])
slide.shapes.title.text = '感谢观看'
box = slide.shapes.add_textbox(Inches(0.7), Inches(2.3), Inches(11.5), Inches(2))
text = box.text_frame
text.text = 'Q&A'
text.paragraphs[0].font.size = Pt(44)

prs.save(str(OUT))
print(OUT)
print(f'images_used={len(images)}')
